#!/usr/bin/env python3
"""
Generate Viet Fire Kitchen proposal as PowerPoint presentation
5 slides, earthy color scheme (greens, browns, cream, gold)
Output: proposals/viet-fire-kitchen.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# Colors: earthy/homey palette
DARK_GREEN = RGBColor(*hex_to_rgb('#3D5A3E'))
CREAM = RGBColor(*hex_to_rgb('#F5F0E8'))
WARM_BROWN = RGBColor(*hex_to_rgb('#8B5E3C'))
GOLD = RGBColor(*hex_to_rgb('#C9A84C'))
TEXT_DARK = RGBColor(40, 40, 40)
WHITE = RGBColor(255, 255, 255)

def add_text_box(slide, left, top, width, height, text, font_size=24, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color

    return textbox

def create_presentation(output_path='proposals/viet-fire-kitchen.pptx'):
    """Create 5-slide Viet Fire Kitchen proposal.

    Args:
        output_path: Path where the presentation will be saved

    Returns:
        Path to the created presentation file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_GREEN

    add_text_box(slide1, Inches(0.5), Inches(2.5), Inches(9), Inches(1.5),
                "Viet Fire Kitchen", font_size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(0.5), Inches(4), Inches(9), Inches(1),
                "Your Marketing System", font_size=40, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(0.5), Inches(5.2), Inches(9), Inches(0.8),
                "Filling Your Slow Hours", font_size=28, color=GOLD, align=PP_ALIGN.CENTER)

    # Slide 2: Problem + What Changes
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide2.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CREAM

    add_text_box(slide2, Inches(0.5), Inches(0.4), Inches(9), Inches(0.6),
                "The Situation", font_size=40, bold=True, color=WARM_BROWN)

    situation = "Packed weekday lunches. Empty dinners. Zero marketing. Understaffed."
    add_text_box(slide2, Inches(0.5), Inches(1.1), Inches(9), Inches(0.8),
                situation, font_size=14, color=TEXT_DARK)

    add_text_box(slide2, Inches(0.5), Inches(2.1), Inches(4.3), Inches(0.4),
                "BEFORE", font_size=16, bold=True, color=WARM_BROWN)
    add_text_box(slide2, Inches(5.2), Inches(2.1), Inches(4.3), Inches(0.4),
                "AFTER", font_size=16, bold=True, color=WARM_BROWN)

    before_items = [
        "No website",
        "No GMB",
        "No offers",
        "No email/SMS",
        "No feedback",
        "No visibility"
    ]

    after_items = [
        "SEO website + ordering",
        "GMB + auto reviews",
        "AI-generated offers",
        "4 emails + 4 SMS/month",
        "Auto follow-up funnel",
        "Targeted campaigns"
    ]

    top = 2.6
    for before, after in zip(before_items, after_items):
        add_text_box(slide2, Inches(0.5), Inches(top), Inches(4.3), Inches(0.35),
                    f"• {before}", font_size=11, color=TEXT_DARK)
        add_text_box(slide2, Inches(5.2), Inches(top), Inches(4.3), Inches(0.35),
                    f"• {after}", font_size=11, color=TEXT_DARK)
        top += 0.42

    # Slide 3: The Solution (5 Systems)
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide3.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CREAM

    add_text_box(slide3, Inches(0.5), Inches(0.4), Inches(9), Inches(0.6),
                "The Solution", font_size=40, bold=True, color=WARM_BROWN)

    systems = [
        "1. Website - SEO-optimized with online ordering link",
        "2. Google My Business - Full optimization + auto review responses",
        "3. Email + SMS Campaigns - 4/month to drive off-peak traffic",
        "4. Offer Generator - AI creates offers, you approve, auto-launches",
        "5. Feedback Loop - Post-order follow-up that drives reviews"
    ]

    top = 1.3
    for system in systems:
        add_text_box(slide3, Inches(0.7), Inches(top), Inches(8.6), Inches(0.45),
                    system, font_size=13, color=TEXT_DARK)
        top += 0.95

    # Slide 4: What We'll Do + Timeline
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide4.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CREAM

    add_text_box(slide4, Inches(0.5), Inches(0.4), Inches(4.3), Inches(0.6),
                "What We'll Do", font_size=28, bold=True, color=WARM_BROWN)
    add_text_box(slide4, Inches(5.2), Inches(0.4), Inches(4.3), Inches(0.6),
                "Timeline", font_size=28, bold=True, color=WARM_BROWN)

    deliverables = [
        "✓ 4 email campaigns/month",
        "✓ 4 SMS blasts/month",
        "✓ AI offer (owner-approved)",
        "✓ 24hr review responses",
        "✓ 1hr post-order follow-up",
        "✓ GMB optimization"
    ]

    top = 1.2
    for item in deliverables:
        add_text_box(slide4, Inches(0.5), Inches(top), Inches(4.3), Inches(0.35),
                    item, font_size=11, color=TEXT_DARK)
        top += 0.4

    add_text_box(slide4, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.45),
                "WEEK 1: Build", font_size=13, bold=True, color=WARM_BROWN)
    add_text_box(slide4, Inches(5.2), Inches(1.7), Inches(4.3), Inches(1.2),
                "• Website build\n• GMB setup\n• Email sequences\n• System test",
                font_size=10, color=TEXT_DARK)

    add_text_box(slide4, Inches(5.2), Inches(3.1), Inches(4.3), Inches(0.45),
                "WEEK 2: Launch", font_size=13, bold=True, color=WARM_BROWN)
    add_text_box(slide4, Inches(5.2), Inches(3.6), Inches(4.3), Inches(1.2),
                "• Final review\n• Customer upload\n• Go live\n• Check-in scheduled",
                font_size=10, color=TEXT_DARK)

    # Slide 5: Investment + Next Steps
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide5.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_GREEN

    add_text_box(slide5, Inches(0.5), Inches(0.4), Inches(9), Inches(0.6),
                "Investment", font_size=40, bold=True, color=GOLD)

    pricing = [
        ("Foundation", "$500 + $200/mo", "Basic"),
        ("Starter", "$800 + $300/mo", "Email only"),
        ("Core (Recommended)", "$1,200 + $400/mo", "FULL SYSTEM"),
        ("Growth", "$1,500 + $500/mo", "+ Auto feedback"),
        ("Partnership", "$800 + $350/mo + 5% revenue", "Aligned incentives")
    ]

    top = 1.3
    for name, price, desc in pricing:
        add_text_box(slide5, Inches(0.5), Inches(top), Inches(9), Inches(0.25),
                    f"{name} - {price}", font_size=11, bold=True, color=WHITE)
        add_text_box(slide5, Inches(0.7), Inches(top+0.28), Inches(8.8), Inches(0.25),
                    f"{desc}", font_size=9, color=GOLD)
        top += 0.65

    add_text_box(slide5, Inches(0.5), Inches(6.2), Inches(9), Inches(1),
                "Next Steps: Choose tier → Schedule kickoff → Launch\n\nQuestions? Text Vu. 48hr response guaranteed.",
                font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Save presentation
    try:
        prs.save(output_path)
        print(f"Proposal created: {output_path}")
        return output_path
    except IOError as e:
        print(f"Error saving presentation: {e}", file=__import__('sys').stderr)
        raise

if __name__ == '__main__':
    import sys
    output_path = sys.argv[1] if len(sys.argv) > 1 else 'proposals/viet-fire-kitchen.pptx'
    create_presentation(output_path)
